import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from flask import Flask, request, send_file
import os
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize

nltk.download('punkt')
nltk.download('stopwords')

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def summarize_text(text):
    stop_words = set(stopwords.words("english"))
    words = word_tokenize(text)
    
    freq_table = dict()
    for word in words:
        word = word.lower()
        if word in stop_words:
            continue
        if word in freq_table:
            freq_table[word] += 1
        else:
            freq_table[word] = 1
    
    sentences = sent_tokenize(text)
    sentence_value = dict()
    
    for sentence in sentences:
        for word, freq in freq_table.items():
            if word in sentence.lower():
                if sentence in sentence_value:
                    sentence_value[sentence] += freq
                else:
                    sentence_value[sentence] = freq
    
    sum_values = 0
    for sentence in sentence_value:
        sum_values += sentence_value[sentence]
    
    average = int(sum_values / len(sentence_value))
    
    summary = ''
    for sentence in sentences:
        if (sentence in sentence_value) and (sentence_value[sentence] > (1.2 * average)):
            summary += " " + sentence
    
    return summary

def split_text_to_slides(text, words_per_slide=200, words_per_line=12):
    words = text.split()
    slides_text = []
    current_slide_words = []

    for word in words:
        current_slide_words.append(word)
        if len(current_slide_words) == words_per_slide:
            slides_text.append(current_slide_words)
            current_slide_words = []

    if current_slide_words:
        slides_text.append(current_slide_words)

    formatted_slides = []
    for slide_words in slides_text:
        slide_text = ""
        for i in range(0, len(slide_words), words_per_line):
            line_words = slide_words[i:i + words_per_line]
            slide_text += " ".join(line_words) + "\n"
        formatted_slides.append(slide_text.strip())

    return formatted_slides

def create_presentation_with_text_and_image(text, image_path):
    presentation = Presentation()
    
    # First slide with image
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    left = top = 0
    pic = slide.shapes.add_picture(image_path, left, top, width=presentation.slide_width, height=presentation.slide_height)
    
    # Second slide with centered text
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    text_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(1), width=Inches(8.5), height=Inches(1))
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = "This is a summarized version from the above book."
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add summarized text to the second slide
    slides_text = split_text_to_slides(text)
    if slides_text:
        slide_text = slides_text.pop(0)
        p = text_frame.add_paragraph()
        p.text = slide_text
        text_frame.word_wrap = True

    # Summarized text in subsequent slides
    for slide_text in slides_text:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        text_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(1), width=Inches(8.5), height=Inches(5))
        text_frame = text_box.text_frame
        text_frame.text = slide_text

    return presentation

def extract_pdf_text_and_first_image(pdf_file_path):
    doc = fitz.open(pdf_file_path)
    full_text = ""
    first_image_path = None

    for page_index, page in enumerate(doc):
        full_text += page.get_text()
        if first_image_path is None:
            images = page.get_images(full=True)
            if images:
                xref = images[0][0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                first_image_path = os.path.join(UPLOAD_FOLDER, "first_image.png")
                with open(first_image_path, "wb") as image_file:
                    image_file.write(image_bytes)

    doc.close()
    return full_text, first_image_path

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return 'No file part', 400
    file = request.files['file']
    if file.filename == '':
        return 'No selected file', 400
    if file:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(file_path)
        extracted_text, first_image_path = extract_pdf_text_and_first_image(file_path)
        summarized_text = summarize_text(extracted_text)
        presentation = create_presentation_with_text_and_image(summarized_text, first_image_path)
        pptx_path = os.path.join(app.config['UPLOAD_FOLDER'], 'presentation.pptx')
        presentation.save(pptx_path)
        return send_file(pptx_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
