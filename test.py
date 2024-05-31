import unittest
from io import BytesIO
from pptx import Presentation
from flask import Flask, request
import fitz  # PyMuPDF
import os

from BookSnap import summarize_text, split_text_to_slides, create_presentation_with_text_and_image, extract_pdf_text_and_first_image

class TestBookSnap(unittest.TestCase):

    def setUp(self):
        self.pdf_path = 'sample.pdf'
        self.upload_folder = 'uploads'
        if not os.path.exists(self.upload_folder):
            os.makedirs(self.upload_folder)
        
        # Create a sample PDF file for testing
        with fitz.open() as pdf:
            page = pdf.new_page()
            page.insert_text((72, 72), "This is a test PDF document. It contains text that will be summarized.")
            pdf.save(self.pdf_path)

    def tearDown(self):
        # Clean up: remove test files and directories
        if os.path.exists(self.pdf_path):
            os.remove(self.pdf_path)
        if os.path.exists(self.upload_folder):
            for file in os.listdir(self.upload_folder):
                os.remove(os.path.join(self.upload_folder, file))
            os.rmdir(self.upload_folder)

    def test_summarize_text(self):
        text = "This is a test PDF document. It contains text that will be summarized. Summarization helps in understanding the content quickly."
        summarized = summarize_text(text)
        self.assertTrue(len(summarized) > 0)
        self.assertFalse("summarized" in summarized)

    def test_split_text_to_slides(self):
        text = " ".join(["word"] * 500)  # create a text with 500 words
        slides = split_text_to_slides(text, words_per_slide=200, words_per_line=12)
        self.assertEqual(len(slides), 3)
        self.assertTrue(all(len(slide.split()) <= 200 for slide in slides))

    def test_create_presentation_with_text_and_image(self):
        text = "This is a summarized text."
        image_path = os.path.join(self.upload_folder, 'test_image.png')
        with open(image_path, 'wb') as f:
            f.write(b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x10\x00\x00\x00\x10\x08\x06\x00\x00\x00\x1f\xf3\xff\xa0\x00\x00\x00\ttIME\x07\xe4\x02\x17\x10\x13\x03%\xf2\xc2\xc0\x00\x00\x00\x19tEXtSoftware\x00Adobe ImageReadyq\xc9e<\x00\x00\x00HIDAT8\x8dcddbf\xa0\x040Q\xa4 \x8b\x81\x81\x81\xe1\x18\x88\x04"\x80\xd8x\x10\xf0\x97\x11H\xc9\xc8\xc8\xc0@\xc0\xe0\xf0\xef\x02\x14\r@0\xe0\x07\n\x0b\n\xe4\x00\x00\x00\x00IEND\xaeB`\x82')
        
        presentation = create_presentation_with_text_and_image(text, image_path)
        pptx_path = os.path.join(self.upload_folder, 'test_presentation.pptx')
        presentation.save(pptx_path)
        
        self.assertTrue(os.path.exists(pptx_path))

        # Check the first slide for the image and text
        presentation = Presentation(pptx_path)
        first_slide = presentation.slides[0]
        self.assertTrue(len(first_slide.shapes) > 0)

        found_image = False
        found_text = False
        for shape in first_slide.shapes:
            if shape.shape_type == 13:  # Picture shape type
                found_image = True
            if shape.has_text_frame and "This is a summarized PowerPoint about the above book." in shape.text:
                found_text = True

        self.assertTrue(found_image)
        self.assertTrue(found_text)

    def test_extract_pdf_text_and_first_image(self):
        text, image_path = extract_pdf_text_and_first_image(self.pdf_path)
        self.assertTrue(len(text) > 0)
        self.assertIsNone(image_path)

if __name__ == '__main__':
    unittest.main()
